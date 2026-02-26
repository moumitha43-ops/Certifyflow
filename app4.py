import streamlit as st
import pandas as pd
from pptx import Presentation
from pathlib import Path
import smtplib
from email.message import EmailMessage
from email.utils import make_msgid
import tempfile
import subprocess
import os

# =========================
# PPT FILL
# =========================
def fill_ppt(template_path, output_ppt, data):
    prs = Presentation(str(template_path))

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    for key, value in data.items():
                        placeholder = f"{{{{{key}}}}}"
                        if placeholder in run.text:
                            run.text = run.text.replace(
                                placeholder,
                                "" if pd.isna(value) else str(value)
                            )

    prs.save(str(output_ppt))


# =========================
# PPT → IMAGE (LIBREOFFICE)
# =========================
import shutil

def ppt_to_images(ppt_path, output_dir):
    output_dir.mkdir(parents=True, exist_ok=True)

    soffice = shutil.which("soffice")

    if not soffice:
        raise RuntimeError(
            "LibreOffice not found. Install LibreOffice or add soffice to PATH."
        )

    cmd = [
        soffice,
        "--headless",
        "--convert-to", "png",
        "--outdir", str(output_dir),
        str(ppt_path),
    ]

    subprocess.run(cmd, check=True)

# =========================
# EMAIL
# =========================
def send_email(sender, password, recipient, subject, name, image_path, event_name):
    msg = EmailMessage()
    msg["From"] = sender
    msg["To"] = recipient
    msg["Subject"] = subject

    cid = make_msgid()

    html = f"""
    <html>
      <body>
        <p>Dear <b>{name}</b>,</p>
        <p>Please find your certificate below:</p>
        <img src="cid:{cid[1:-1]}" width="700">
        <p>Regards,<br>Event Team</p>
      </body>
    </html>
    """

    msg.set_content("Your certificate is attached.")
    msg.add_alternative(html, subtype="html")

    with open(image_path, "rb") as img:
        msg.get_payload()[1].add_related(
            img.read(),
            maintype="image",
            subtype="png",
            cid=cid,
            filename=f"{event_name}.png",
            disposition="inline"
        )

    with smtplib.SMTP_SSL("smtp.gmail.com", 465) as server:
        server.login(sender, password)
        server.send_message(msg)


# =========================
# STREAMLIT UI
# =========================
st.set_page_config("Certificate Generator", layout="centered")
st.title("🎓 Cloud Certificate Generator")

event_name = st.text_input("📌 Event Name")
ppt_file = st.file_uploader("📄 PPT Template", type="pptx")
csv_file = st.file_uploader("📊 CSV File", type="csv")

st.subheader("📧 Email Settings")
sender_email = st.text_input("Sender Email")
sender_password = st.text_input("Gmail App Password", type="password")
email_column = st.text_input("Email column name", value="EMAIL")

# =========================
# PREVIEW
# =========================
if csv_file and ppt_file:
    df = pd.read_csv(csv_file)

    if "NAME" in df.columns:
        name_preview = st.selectbox("Preview Name", df["NAME"].astype(str))

        if st.button("👁 Preview Certificate"):
            row = df[df["NAME"].astype(str) == name_preview].iloc[0].to_dict()

            with tempfile.TemporaryDirectory() as tmp:
                template = Path(tmp) / "template.pptx"
                template.write_bytes(ppt_file.read())

                ppt_out = Path(tmp) / "preview.pptx"
                img_out = Path(tmp) / "img"

                fill_ppt(template, ppt_out, row)
                ppt_to_images(ppt_out, img_out)

                images = list(img_out.glob("*.png"))
                if images:
                    st.image(str(images[0]), use_column_width=True)

# =========================
# GENERATE & SEND
# =========================
if st.button("🚀 Generate & Send Certificates"):
    if not all([event_name, ppt_file, csv_file, sender_email, sender_password]):
        st.error("Fill all fields")
        st.stop()

    df = pd.read_csv(csv_file)
    progress = st.progress(0)

    with tempfile.TemporaryDirectory() as tmp:
        template_path = Path(tmp) / "template.pptx"
        template_path.write_bytes(ppt_file.read())

        for idx, row in df.iterrows():
            data = row.to_dict()
            name = str(data.get("NAME", f"user_{idx}"))
            recipient = data.get(email_column)

            ppt_out = Path(tmp) / f"{name}.pptx"
            img_out = Path(tmp) / name

            fill_ppt(template_path, ppt_out, data)
            ppt_to_images(ppt_out, img_out)

            images = list(img_out.glob("*.png"))
            if recipient and images:
                send_email(
                    sender_email,
                    sender_password,
                    recipient,
                    f"{event_name} Certificate",
                    name,
                    images[0],
                    event_name
                )

            progress.progress(int((idx + 1) / len(df) * 100))

    st.success("🎉 All certificates sent successfully!")